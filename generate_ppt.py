from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs):
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(28, 25, 23) # Stone 900
    
    # Accent bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(1), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(185, 28, 28) # Red 700
    shape.line.fill.background()
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "IATF 16949 智能平台"
    p.font.bold = True
    p.font.size = Pt(54)
    p.font.color.rgb = RGBColor(245, 245, 244) # Stone 100
    p.font.name = "Microsoft YaHei"
    
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "数字化时代的汽车行业质量标准桌面工作台"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(168, 162, 158) # Stone 400
    p2.font.name = "Microsoft YaHei"

def add_content_slide(prs, title_text, bullet_points, screenshot_hint):
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Light Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 249) # Stone 50
    
    # Top Header Bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(28, 25, 23) # Stone 900
    shape.line.fill.background()
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "Microsoft YaHei"
    
    # Left Content (Bullets)
    txBox_bullets = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(4))
    tf_bullets = txBox_bullets.text_frame
    tf_bullets.word_wrap = True
    for bp in bullet_points:
        p = tf_bullets.add_paragraph()
        p.text = "• " + bp
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(41, 37, 36) # Stone 800
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(14)
        
    # Right Content (Screenshot Placeholder)
    placeholder = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(1.8), Inches(4.5), Inches(4.5))
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RGBColor(231, 229, 228) # Stone 200
    placeholder.line.color.rgb = RGBColor(168, 162, 158)
    
    # Placeholder Text
    txBox_ph = slide.shapes.add_textbox(Inches(5), Inches(3.5), Inches(4.5), Inches(1))
    tf_ph = txBox_ph.text_frame
    tf_ph.word_wrap = True
    p_ph = tf_ph.paragraphs[0]
    p_ph.text = screenshot_hint
    p_ph.alignment = PP_ALIGN.CENTER
    p_ph.font.bold = True
    p_ph.font.size = Pt(18)
    p_ph.font.color.rgb = RGBColor(185, 28, 28) # Red 700
    p_ph.font.name = "Microsoft YaHei"

def add_closing_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(185, 28, 28) # Red 700
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "赋能质量人，提升企业体系能力"
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "Microsoft YaHei"
    
    # QR Placeholder
    qr_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.5), Inches(2), Inches(2))
    qr_shape.fill.solid()
    qr_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    qr_shape.line.fill.background()
    
    txBox_qr = slide.shapes.add_textbox(Inches(4), Inches(5.2), Inches(2), Inches(1))
    p_qr = txBox_qr.text_frame.paragraphs[0]
    p_qr.text = "贴入网站二维码"
    p_qr.alignment = PP_ALIGN.CENTER
    p_qr.font.size = Pt(14)
    p_qr.font.color.rgb = RGBColor(185, 28, 28)
    p_qr.font.name = "Microsoft YaHei"

def create_presentation():
    prs = Presentation()
    
    # 设置幻灯片比例为宽屏 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 1. 封面
    add_title_slide(prs)
    
    # 2. 痛点
    add_content_slide(
        prs,
        "传统标准学习的痛点",
        [
            "标准文本极其枯燥，纯理论难落地",
            "遇到审核问题时，PDF翻找效率低下",
            "只记住了编号，却不知道在工厂怎么留存证据"
        ],
        "【请在此贴上一张厚重且密密麻麻的标准PDF截图】\n体现出传统方式的压抑感"
    )
    
    # 3. 浏览功能
    add_content_slide(
        prs,
        "特性一：结构化浏览与毫秒级速查",
        [
            "打破几百页纸的限制，将条款转化为树状层级数据",
            "支持按编号和关键字进行全局毫秒级搜索",
            "护眼的暗黑模式与流畅的交互动画，缓解视觉疲劳"
        ],
        "【请在此贴上 BrowseMode (浏览模式) 的深色界面截图】"
    )

    # 4. 解读功能
    add_content_slide(
        prs,
        "特性二：白话解析与落地证据指南",
        [
            "内嵌独家条款解读，把抽象的行话翻译成车间术语",
            "一键提取“核心要点”，直击内审员的审查重点",
            "明确每个条款的实际业务应用场景和输出文档要求"
        ],
        "【请在此贴上条款展开后，包含红色‘实际应用’标签的页面截图】"
    )
    
    # 5. 测验功能
    add_content_slide(
        prs,
        "特性三：场景化学习与在线测验",
        [
            "采用抽认卡（Flashcards）进行沉浸式记忆训练",
            "内置模拟测验模式，即时反馈，找出知识盲区",
            "直观的进度环和正确率统计，让学习成就感拉满"
        ],
        "【请在此贴上 QuizMode (在线测验) 或 抽认卡 翻转效果的截图】"
    )
    
    # 6. 结尾
    add_closing_slide(prs)
    
    prs.save('IATF16949_Smart_Platform_Pro.pptx')

if __name__ == '__main__':
    create_presentation()
    print("Beautiful PPT generated successfully!")
